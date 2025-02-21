import re
from tqdm import tqdm
import fitz
from langchain_core.documents import Document
from pptx import Presentation
from docx import Document as DocxDocument
import openpyxl
import csv
from langchain_community.vectorstores import Milvus
from langchain_huggingface import HuggingFaceEmbeddings
from .enable_logging import logger 
from cohere_app.Chunking_UI import db_utility
from doctr.io import DocumentFile
from doctr.models import ocr_predictor

embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2', model_kwargs={'device': "cuda"})
model = ocr_predictor(pretrained=True)
model = model.to('cuda')
OCR_LIST = []

def extract_text_pdf(pdf_path):
    text_by_page = []
    try:
        pdf = fitz.open(pdf_path)
    except Exception as e:
        return [], "Can't open the file ERROR"
    for page_num in range(len(pdf)):
        page = pdf[page_num]
        text = page.get_text()
        # we need to clean the text here.
        if not text:
            OCR_LIST.append(pdf_path)
            return [], "OCR required - ERROR"
        else:
            text_by_page.append((page_num, text))
    return text_by_page, "Text extraction done"

def process_ocr_document(pdf_path):
    """ Extract text using OCR from each page of the PDF """
    try:
        """ Extract text using OCR from each page of the PDF """
        text_by_page = []
        doc = DocumentFile.from_pdf(pdf_path)
        result = model(doc)
        for page_num, page in enumerate(result.pages, start=1):
            page_text = ""
            words = []
            for block in page.blocks:
                for line in block.lines:
                    for word in line.words:
                        words.append(word)
            sorted_words = sorted(words, key=lambda word: (word.geometry[0][1], word.geometry[0][0]))
            page_text = " ".join([word.value for word in sorted_words])
            text_by_page.append((page_num, page_text))
        return text_by_page, "text extraction done"
    except Exception as e:
        logger.error(f"Error extracting OCR text from {pdf_path}: {e}")
        return [], f"OCR FILE ERROR : {e}"


def process_pptx(file_path):
    """ Process .pptx file and extract text slide by slide """
    try:
        prs = Presentation(file_path)
        text_by_slide = []
        for i, slide in enumerate(prs.slides):
            slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            text_by_slide.append((i + 1, slide_text))
        return text_by_slide, "text extraction done" 
    except Exception as e:
        logger.exception(f"Error processing PPTX file {file_path}")
        return [], f"PPTX FILE ERROR : {e}"


def process_docx(file_path):
    """ Process .docx file and extract text """
    try:
        doc = DocxDocument(file_path)
        text_by_paragraph = [(i + 1, para.text) for i, para in enumerate(doc.paragraphs)]
        return text_by_paragraph, "text extraction done" 
    except Exception as e:
        logger.exception(f"Error processing DOCX file {file_path}")
        return [], f"DOCX FILE ERROR : {e}"


def process_txt(file_path):
    """ Process .txt file and extract text """
    try:
        with open(file_path, 'r') as f:
            text = f.read()
        return [(1, text)], "text extraction done"  # Text file considered as one page
    except Exception as e:
        logger.error(f"Error processing TXT file {file_path}: {e}")
        return [], f"TEXT FILE ERROR : {e}"


def process_xlsx(file_path):
    """ Process .xlsx file and extract text sheet by sheet """
    try:
        workbook = openpyxl.load_workbook(file_path)
        text_by_sheet = []
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            sheet_text = "\n".join([",".join([str(cell.value) for cell in row]) for row in worksheet.iter_rows()])
            text_by_sheet.append((sheet, sheet_text))
        return text_by_sheet, "text extraction done"
    except Exception as e:
        logger.exception(f"Error processing XLSX file {file_path}")
        return [], f"EXCEL FILE ERROR : {e}"


def process_csv(file_path):
    """ Process .csv file and extract text """
    try:
        with open(file_path, 'r') as f:
            reader = csv.reader(f)
            text = "\n".join([",".join(row) for row in reader])
        return [(1, text)], "text extraction done"  # Consider CSV as one page
    except Exception as e:
        logger.error(f"Error processing CSV file {file_path}: {e}")
        return [()],  f"CSV FILE ERROR : {e}"


def clean_text(text):
    """ Clean the extracted text by removing unwanted characters and formatting """
    text = re.sub(r'\.{3,}', '.', text)
    lines = text.split('\n')
    cleaned_lines = [line.strip() for line in lines if len(line.split()) >= 4]
    cleaned_text = '\n'.join(cleaned_lines)
    cleaned_text = re.sub(r'^\s*$', '', cleaned_text, flags=re.MULTILINE)
    return cleaned_text.strip()

def clean_chunk(chunk):
    """
    Clean a text chunk by removing extra formatting characters.
    Processes text line-by-line, removing lines that consist mostly of formatting symbols,
    then collapses extra whitespace.
    """
    lines = chunk.splitlines()
    cleaned_lines = []
    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue
        alnum_count = sum(1 for ch in stripped_line if ch.isalnum())
        if len(stripped_line) > 0 and (alnum_count / len(stripped_line)) < 0.3:
            continue
        cleaned_lines.append(line)
    cleaned_chunk = " ".join(cleaned_lines)
    cleaned_chunk = re.sub(r'[\-\+\|=]{2,}', ' ', cleaned_chunk)
    cleaned_chunk = re.sub(r'\s+', ' ', cleaned_chunk)
    return cleaned_chunk.strip()

def read_and_split_text(text_by_page, chunk_size=800, overlap_size=200):
    """
    Create chunks using a sliding window approach while tracking page numbers.
    Each chunk is built by concatenating cleaned paragraphs until it reaches at least
    chunk_size characters, then extended until the end of the sentence if needed.
    Returns a list of tuples: (chunk_text, page_number).
    """
    logger.info(f"Creating chunks from {len(text_by_page)} pages...")
    
    chunks = []
    current_text = ""
    current_pages = set()
    
    for page_num, text in text_by_page:
        cleaned_para = clean_chunk(text)
        if not cleaned_para:
            continue
            
        current_text += cleaned_para + " "
        current_pages.add(page_num)
        
        while len(current_text) >= chunk_size:
            chunk = current_text[:chunk_size]
            
            if len(current_text) > chunk_size and current_text[chunk_size] == '.':
                chunk = current_text[:chunk_size+1]
            elif not chunk.rstrip().endswith('.'):
                period_index = current_text.find('.', chunk_size)
                if period_index != -1:
                    chunk = current_text[:period_index + 1]
            
            chunk = chunk.strip()
            if chunk:
                start_page = min(current_pages) if current_pages else page_num
                end_page = max(current_pages) if current_pages else page_num
                chunks.append((chunk, start_page, end_page))
            
            current_text = current_text[len(chunk) - overlap_size:].strip()
            current_pages = {page_num}
    
    if current_text.strip():
        start_page = min(current_pages) if current_pages else page_num
        end_page = max(current_pages) if current_pages else page_num
        chunks.append((current_text.strip(), start_page, end_page))
    
    logger.info(f"Created {len(chunks)} chunks")
    return chunks


def process_document(file_path):
    """ Process a single document """
    try:
        text_by_page = []
        message = ""
        if file_path.endswith('.pdf') or file_path.endswith('.PDF'):
            text_by_page, message = extract_text_pdf(file_path)
        elif file_path.endswith('.pptx'):
            text_by_page, message = process_pptx(file_path)
        elif file_path.endswith('.docx'):
            text_by_page, message = process_docx(file_path)
        elif file_path.endswith('.txt'):
            text_by_page, message = process_txt(file_path)
        elif file_path.endswith('.xlsx'):
            text_by_page, message = process_xlsx(file_path)
        elif file_path.endswith('.csv'):
            text_by_page, message = process_csv(file_path)
        return text_by_page, message
    #TODO - i dont think so this except requires to store the error files
    except Exception as e:
        return [], f"Error processing document {file_path}\n{e}"

def create_langchain_documents(found_files, collection_name):
    """
    Create langchain documents from the extracted and processed text.
    Now uses sliding window chunking for better text segmentation.
    """
    db_utility.create_user_access(collection_name)
    db_utility.chunking_monitor()
    db_utility.create_error_files(collection_name)

    # Process found files
    for file_index, file in enumerate(found_files):
        text_by_page, message = process_document(file)
        if text_by_page and 'error' not in message.lower():
            chunks = read_and_split_text(text_by_page)
            logger.info(f"Current processing file {file} with {len(chunks)} chunks")
            
            if chunks:
                documents = []
                for chunk, start_page, end_page in chunks:
                    doc = Document(
                        page_content=chunk,
                        metadata={
                            'source': file,
                            'page': f"{start_page}-{end_page}" if start_page != end_page else str(start_page)
                        }
                    )
                    documents.append(doc)

                if documents:
                    try:
                        Milvus.from_documents(
                            documents,
                            embeddings,
                            collection_name=collection_name,
                            connection_args={'uri': "http://localhost:19530"}
                        )
                        db_utility.insert_user_access(file, 'YES', message, collection_name)
                    except Exception as e:
                        error_message = f"Error inserting into Milvus: {str(e)}"
                        logger.error(error_message)
                        db_utility.store_error_files_with_error(collection_name, file, error_message)
            else:
                db_utility.store_error_files_with_error(collection_name, file, "No valid chunks generated")
                
        elif "error" in message.lower() and "ocr" not in message.lower():
            db_utility.store_error_files_with_error(collection_name, file, message)
            logger.error(f"Error in the document - Skipping {file}")
            
        progress_percentage = (file_index + 1) / len(found_files) * 100
        yield {"progress_percentage": progress_percentage, "current_progress": file_index + 1, "total_files": len(found_files)}

    for ocr_file_index, ocr_file in enumerate(tqdm(OCR_LIST, desc="Overall Progress")):
        tqdm().set_description(f"Processing File: {ocr_file}")
        text_by_page, message = process_ocr_document(ocr_file)
        if text_by_page:
            chunks = read_and_split_text(text_by_page)
            logger.info(f"Current processing OCR file {ocr_file} with {len(chunks)} chunks")
            
            if chunks:
                documents = []
                for chunk, start_page, end_page in chunks:
                    doc = Document(
                        page_content=chunk,
                        metadata={
                            'source': ocr_file,
                            'page': f"{start_page}-{end_page}" if start_page != end_page else str(start_page)
                        }
                    )
                    documents.append(doc)
                    
                if documents:
                    try:
                        Milvus.from_documents(
                            documents,
                            embeddings,
                            collection_name=collection_name,
                            connection_args={'uri': "http://localhost:19530"}
                        )
                        db_utility.update_ocr_status(ocr_file, collection_name)
                    except Exception as e:
                        error_message = f"Error inserting OCR document into Milvus: {str(e)}"
                        logger.error(error_message)
                        db_utility.store_error_files_with_error(collection_name, ocr_file, error_message)
                        
        progress_percentage = (ocr_file_index + 1) / len(OCR_LIST) * 100
        yield {"progress_percentage": progress_percentage, "current_progress": ocr_file_index + 1, "total_files": len(OCR_LIST)}
